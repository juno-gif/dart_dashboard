"""
PPT 생성 서비스 — Story 6.1
분석 세트의 재무 데이터를 PowerPoint 파일로 변환
[Source: architecture.md - Service Layer Patterns]
"""
import io
from datetime import date

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def generate_analysis_ppt(
    set_name: str,
    company_codes: list[str],
    financials_by_corp: dict,
) -> bytes:
    """분석 세트 PPT 생성.

    Args:
        set_name: 분석 세트 이름 (제목 슬라이드에 사용)
        company_codes: 기업 코드 목록 (슬라이드 순서 결정)
        financials_by_corp: {corp_code: [재무 데이터 row dict, ...]} 형태

    Returns:
        PPT 파일 bytes
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: 제목 슬라이드
    _add_title_slide(prs, set_name, company_codes)

    # Slide 2~N: 기업별 P&L 트렌드
    for corp_code in company_codes:
        data = financials_by_corp.get(corp_code, [])
        _add_company_pl_slide(prs, corp_code, data)

    # Slide N+1: 전체 기업 매출 비교
    _add_comparison_slide(prs, company_codes, financials_by_corp)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(prs: Presentation, set_name: str, company_codes: list[str]) -> None:
    """Slide 1: 제목 슬라이드"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = set_name
    slide.placeholders[1].text = (
        f"생성일: {date.today().isoformat()}\n기업: {', '.join(company_codes)}"
    )


def _add_company_pl_slide(prs: Presentation, corp_code: str, financials: list) -> None:
    """기업별 P&L 트렌드 슬라이드 (매출·영업이익·순이익 바 차트)"""
    # 연도별 계정과목 그룹핑 (억 단위 변환)
    by_year: dict = {}
    for row in financials:
        year = row["bsns_year"]
        account_key = row["account_key"]
        amount = (row.get("amount") or 0) / 100_000_000
        by_year.setdefault(year, {})[account_key] = amount
    years = sorted(by_year.keys())

    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 제목 텍스트박스
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    tf.text = f"{corp_code} P&L 트렌드 (억원)"
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.bold = True

    if not years:
        # 데이터 없음 안내
        nodata_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1))
        nodata_box.text_frame.text = "재무 데이터가 없습니다."
        return

    chart_data = ChartData()
    chart_data.categories = years
    chart_data.add_series("매출", [by_year[y].get("revenue", 0) for y in years])
    chart_data.add_series("영업이익", [by_year[y].get("operating_profit", 0) for y in years])
    chart_data.add_series("순이익", [by_year[y].get("net_income", 0) for y in years])

    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5),
        Inches(1.0),
        Inches(12),
        Inches(6),
        chart_data,
    )


def _add_comparison_slide(
    prs: Presentation,
    company_codes: list[str],
    financials_by_corp: dict,
) -> None:
    """전체 기업 매출 비교 슬라이드 (가장 최근 연도 기준)"""
    latest_revenues: dict = {}
    for corp_code in company_codes:
        data = financials_by_corp.get(corp_code, [])
        by_year: dict = {}
        for row in data:
            year = row["bsns_year"]
            account_key = row["account_key"]
            amount = (row.get("amount") or 0) / 100_000_000
            by_year.setdefault(year, {})[account_key] = amount
        if by_year:
            latest_year = sorted(by_year.keys())[-1]
            latest_revenues[corp_code] = by_year[latest_year].get("revenue", 0)

    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "기업별 매출 비교 (억원, 최근 연도)"
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.bold = True

    if not latest_revenues:
        nodata_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1))
        nodata_box.text_frame.text = "비교할 재무 데이터가 없습니다."
        return

    chart_data = ChartData()
    chart_data.categories = list(latest_revenues.keys())
    chart_data.add_series("매출", list(latest_revenues.values()))

    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5),
        Inches(1.0),
        Inches(12),
        Inches(6),
        chart_data,
    )
